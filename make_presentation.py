import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

def create_extended_presentation():
    prs = Presentation()
    
    # --- HELPER: Add Slide ---
    def add_slide(title_text, content_points=None, layout_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        title = slide.shapes.title
        title.text = title_text
        
        # Formatting title
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

        if content_points:
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = content_points[0] # First point
            
            for point in content_points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
            # Format content
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.space_after = Pt(10)
        
        return slide

    # --- SLIDE 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Further development of an AI-Supported Algorithm\nto choose a suitable Shaft-Hub-Connection"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = ("Bachelor Thesis Defense\n\n"
                     "Author: Farah Hany\n"
                     "Supervisor: Prof. Dr.-Ing. Michael Lätzer\n"
                     "Technische Hochschule Ulm | Faculty of Computer Science\n"
                     "January 2026")
    
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Good morning/afternoon everyone. My name is Farah Hany.\n"
                  "Today I am presenting my Bachelor Thesis on the development of an AI-supported algorithm for selecting shaft-hub connections.\n"
                  "This project focuses on bridging the gap between mechanical engineering standards and modern machine learning techniques.")

    # --- SLIDE 2: Introduction ---
    content = [
        "The Shaft-Hub-Connection (SHC) Challenge:",
        "• Fundamental component in mechanical drive trains.",
        "• Function: Transmit torque from a shaft to a hub (gear, pulley).",
        "• Critical Design Trade-offs:",
        "   - Torque Capacity (Safety)",
        "   - Assembly/Disassembly requirements",
        "   - Manufacturing Costs",
        "• Current Practice: Manual calculation using DIN handbooks."
    ]
    slide = add_slide("Introduction: The Selection Challenge", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Shaft-hub connections are everywhere—in cars, wind turbines, and robotics.\n"
                  "The challenge is that selecting the 'right' connection is a complex trade-off.\n"
                  "A press fit is cheap but hard to assemble. A spline is strong but expensive.\n"
                  "Currently, engineers calculate this manually using heavy standards like DIN 7190. It's time-consuming and prone to human error.")

    # --- SLIDE 3: Problem Statement ---
    content = [
        "Why AI? Why Now?",
        "• Goal: Automate the decision process for instant, safe recommendations.",
        "The Core Problem (Research Gap):",
        "• Machine Learning requires data.",
        "• There are NO publicly available, labeled datasets for SHC selection.",
        "• Historical company data is often fragmented or non-existent.",
        "• Therefore: A standard supervised learning approach was impossible."
    ]
    slide = add_slide("Problem Statement & Research Gap", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("We wanted to automate this with AI.\n"
                  "But we hit a major wall: Data.\n"
                  "Unlike image recognition where you have ImageNet, there is no 'ShaftNet' for engineering.\n"
                  "Without historical data, you cannot train a model. This was the primary problem my thesis had to solve.")

    # --- SLIDE 4: Objectives ---
    content = [
        "1. Analytical Engine: Develop a physics-based calculator (DIN Standards).",
        "2. Synthetic Data: Use the analytical engine to generate a training dataset.",
        "3. Machine Learning: Train and evaluate models (Random Forest, CatBoost).",
        "4. Application: Build a web-based tool for real-time usage.",
        "5. Transparency: Ensure results are explainable (Hybrid output)."
    ]
    slide = add_slide("Research Objectives", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("To solve this, I defined 5 clear objectives:\n"
                  "1. First, build a robust analytical engine based on physics.\n"
                  "2. Use that engine to generate synthetic data—essentially teaching the AI the laws of physics.\n"
                  "3. Train the model.\n"
                  "4. Build a usable web app.\n"
                  "5. Ensure transparency so engineers trust the result.")

    # --- SLIDE 5: Background (Connection Types) ---
    content = [
        "Three Connection Types Analyzed:",
        "1. Interference (Press) Fit (DIN 7190):",
        "   - Friction closure. Permanent. High concentricity.",
        "2. Parallel Key (DIN 6885):",
        "   - Form closure. Easy assembly. Weak under reversing loads.",
        "3. Spline (DIN 5480):",
        "   - Form closure. High torque capacity. Expensive manufacturing.",
        "Scope: Standardized logic for these three most common types."
    ]
    slide = add_slide("Background: Connection Standards", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("I focused on the three industry standards.\n"
                  "Press fits rely on friction—great for high speed, bad for disassembly.\n"
                  "Keys are the standard 'cheap and easy' option but can fail under shock loads.\n"
                  "Splines are the heavy-duty option—complex teeth, very strong, very expensive.")

    # --- SLIDE 6: Methodology Overview ---
    slide = add_slide("Methodology: The 4-Stage Pipeline")
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(4.5)
    shape = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "[INSERT FIGURE 2.16 (System Architecture)]\n\nUser Input -> Analytical Calculation -> \nSynthetic Data Generation -> ML Training -> Final Hybrid App"
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("This is the high-level architecture of my work.\n"
                  "It starts with the Analytical Selector.\n"
                  "This feeds into the Data Generator.\n"
                  "The data trains the Model.\n"
                  "And finally, everything is wrapped in a Hybrid Application that runs both the math and the AI in parallel.")

    # --- SLIDE 7: Methodology - Analytical Engine ---
    content = [
        "The Logic Core (The 'Oracle'):",
        "• Calculates maximum torque capacity for a given geometry.",
        "• Checks Manufacturability Constraints:",
        "   - Press Fit: Is interference < 0.02mm? (Assembly check)",
        "   - Key: Is shaft diameter > minimum required?",
        "• Preference Scoring:",
        "   - Assigns scores based on user priorities (Cost, Assembly, Weight).",
        "   - Resolves 'ties' where multiple connections are technically feasible."
    ]
    slide = add_slide("Stage 1: The Analytical Engine", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("The first stage was coding the DIN standards into Python.\n"
                  "This engine doesn't just check torque. It checks manufacturability.\n"
                  "For example, a press fit might be strong enough, but if the interference is too high, you can't assemble it.\n"
                  "It also applies a scoring logic to decide the 'best' option when more than one type works.")

    # --- SLIDE 8: Methodology - Synthetic Data ---
    content = [
        "Generating the 'ShaftNet' Dataset:",
        "• Input Space Sampling: Randomly sampled parameters within DIN ranges.",
        "   - Shaft Diameter: 6mm - 230mm",
        "   - Torque: 10 Nm - 10,000 Nm",
        "   - Materials: E295, C45E, 42CrMo4, etc.",
        "• Labeling: The Analytical Engine labeled each sample.",
        "• Filtering: Infeasible designs (e.g., plastic deformation) were removed.",
        "• Result: 4,993 Clean, Labeled Samples."
    ]
    slide = add_slide("Stage 2: Synthetic Data Generation", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("This is the most innovative part of the thesis.\n"
                  "I wrote a script to generate 5,000 random engineering scenarios.\n"
                  "My analytical engine acted as the 'Oracle', calculating the correct answer for each.\n"
                  "This gave me a clean, balanced dataset covering the entire range of standard shaft sizes, from 6mm to 230mm.")

    # --- SLIDE 9: Methodology - ML Models ---
    content = [
        "Model Selection:",
        "• Task: Multi-class Classification (Output: Press Fit, Key, or Spline).",
        "• Models Evaluated:",
        "   1. Random Forest (Baseline)",
        "   2. XGBoost (Gradient Boosting)",
        "   3. CatBoost (Categorical Boosting)",
        "• Why CatBoost Won?",
        "   - Native handling of categorical features (Materials).",
        "   - Best performance on tabular data without extensive tuning.",
        "   - Fast inference time (<50ms)."
    ]
    slide = add_slide("Stage 3: Machine Learning Strategy", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("For the AI, I compared three tree-based models.\n"
                  "Deep learning was unnecessary for this tabular data.\n"
                  "CatBoost was the winner. It handles categorical variables—like Steel types—natively, without needing complex encoding.\n"
                  "It was also the fastest, which is important for the web app.")

    # --- SLIDE 10: Results - Analytical Verification ---
    slide = add_slide("Results: Analytical Verification")
    # Table logic same as before but emphasized
    rows, cols = 4, 4
    left, top, width, height = Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["Metric", "Press Fit", "Parallel Key", "Spline"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
    data = [
        ["Torque Cap.", "3,951 Nm", "945 Nm", "6,470 Nm"],
        ["% of Design", "227%", "54.3% (FAIL)", "371.8%"],
        ["Result", "Infeasible*", "Infeasible", "FEASIBLE"]
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(18)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    p = textbox.text_frame.add_paragraph()
    p.text = "Test Case: 45mm Shaft, 870 Nm Required.\n*Press Fit failed due to high interference (>0.02mm)."
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Before looking at AI scores, I had to verify the physics.\n"
                  "I ran a standard test case: 45mm shaft, 870 Nm.\n"
                  "The Key failed—too weak.\n"
                  "The Press fit was strong enough but failed the 'manufacturability' check.\n"
                  "Only the Spline passed.\n"
                  "This matched the manual DIN calculation perfectly.")

    # --- SLIDE 11: Results - Dataset Analysis ---
    content = [
        "Dataset Characteristics:",
        "• Class Balance:",
        "   - Splines: ~38% (High torque preference)",
        "   - Keys: ~32% (Low cost preference)",
        "   - Press Fits: ~30%",
        "• Correlation Insights:",
        "   - Strong correlation between Shaft Diameter and Torque Capacity.",
        "   - 'Cost Preference' strongly negatively correlated with Splines.",
        "• This confirms the dataset represents real-world engineering physics."
    ]
    slide = add_slide("Results: Dataset Insights", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Analyzing the generated data itself was interesting.\n"
                  "We see a roughly even split, but Splines are slightly more common because they handle the high-torque cases where others fail.\n"
                  "The correlations confirmed that the 'Oracle' was working: as Cost preference goes up, Spline recommendations go down.")

    # --- SLIDE 12: Results - Model Performance ---
    slide = add_slide("Results: ML Metrics (CatBoost)")
    chart_data = CategoryChartData()
    chart_data.categories = ['Press Fit', 'Parallel Key', 'Spline', 'Macro Avg']
    chart_data.add_series('F1-Score', (0.6774, 0.8057, 0.9127, 0.7986))
    x, y, w, h = Inches(2), Inches(2), Inches(6), Inches(4)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 1.0
    chart.chart_title.text_frame.text = "F1-Score by Class"
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Here are the AI results.\n"
                  "Overall Macro F1 score is roughly 0.80.\n"
                  "Splines (0.91) are very easy to distinguish—they are the distinct 'heavy duty' choice.\n"
                  "Press Fits (0.68) are the hardest. They often overlap with Keys in the 'medium' range, confusing the model.")

    # --- SLIDE 13: Results - Confusion Matrix ---
    slide = add_slide("Results: Confusion Matrix")
    left = Inches(2.0)
    top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(5.0)
    shape = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "[INSERT CONFUSION MATRIX IMAGE]\n\nHighlight the confusion between Press Fit and Parallel Key."
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("This Confusion Matrix details the errors.\n"
                  "You can see the diagonal is strong—most predictions are correct.\n"
                  "However, notice the block between Press Fit and Key.\n"
                  "In 'borderline' cases where both work, the AI sometimes struggles to guess which one the preference slider favored.")

    # --- SLIDE 14: Results - Feature Importance ---
    slide = add_slide("Results: Feature Importance")
    left = Inches(2.0)
    top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(5.0)
    shape = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "[INSERT FEATURE IMPORTANCE PLOT]\n\nShows: Diameter, Torque, and Cost Preference as top features."
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("What does the AI actually look at?\n"
                  "As expected, Shaft Diameter and Torque are the top drivers.\n"
                  "But interestingly, 'Cost Preference' is also very high.\n"
                  "This proves the model successfully learned the user's subjective preferences, not just the objective physics.")

    # --- SLIDE 15: Web Application - Stack ---
    content = [
        "Tech Stack:",
        "• Backend: FastAPI (Python)",
        "   - Serves the CatBoost model.",
        "   - Runs the Analytical Verification script.",
        "• Frontend: React.js",
        "   - Dynamic form for Geometry inputs.",
        "   - Sliders for User Preferences (0-10 scale).",
        "• Deployment: Localhost / Docker container capability."
    ]
    slide = add_slide("Web Application Architecture", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("To make this usable, I built a web app.\n"
                  "The backend uses FastAPI because it's fast and python-native, easy to integrate with ML libraries.\n"
                  "The frontend is React, providing a modern, responsive interface.\n"
                  "It takes inputs, sends them to the backend, and receives both the ML prediction and the physics check.")

    # --- SLIDE 16: Web Application - Demo ---
    slide = add_slide("Web App Demonstration")
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.0)
    shape = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "[INSERT SCREENSHOT OF WEB APP]\n\nShow: Input fields on left, Results panel on right."
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Here is the final application.\n"
                  "On the left, the engineer enters the shaft diameter, torque, and material.\n"
                  "They also set their preferences: Do they care more about Cost? Or Assembly?\n"
                  "On the right, they get an instant recommendation.")

    # --- SLIDE 17: Discussion ---
    content = [
        "Interpreting the Results:",
        "• Success: The model learned to approximate DIN standards with 80% accuracy.",
        "• Speed vs. Accuracy: ML is instant; Analytical takes ~200ms.",
        "   - In a web app, both are fast enough.",
        "   - In a large-scale optimization (1M+ parts), ML would be critical.",
        "• The 'Hybrid' Safety Net:",
        "   - ML gives the 'suggestion'.",
        "   - Analytical Engine gives the 'proof'.",
        "   - This eliminates the risk of AI 'hallucination' in safety-critical tasks."
    ]
    slide = add_slide("Discussion", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("So, what does this mean?\n"
                  "We proved AI can learn engineering standards.\n"
                  "But more importantly, we showed that a Hybrid approach is the safe path forward.\n"
                  "We don't trust the AI blindly. We use it to guide us, but we let the analytical engine verify the safety factors.\n"
                  "This dual-check system is essential for safety-critical mechanical engineering.")

    # --- SLIDE 18: Limitations ---
    content = [
        "Limitations of the Current System:",
        "• Static Loading Only:",
        "   - Does not account for fatigue, shock, or reversing loads.",
        "• Limited Scope:",
        "   - Only 3 connection types (No polygonal profiles, clamping sets).",
        "   - Only standard DIN materials.",
        "• Synthetic Data Bias:",
        "   - The model is only as good as the 'Oracle' logic.",
        "   - Any simplification in the analytical code is learned by the AI."
    ]
    slide = add_slide("Limitations", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Of course, there are limitations.\n"
                  "I only considered static loads. In the real world, fatigue is a huge factor.\n"
                  "I also limited it to the three most common connections.\n"
                  "And finally, the 'Oracle' bias: if my analytical math has a simplification, the AI learns that simplification as truth.")

    # --- SLIDE 19: Future Work ---
    content = [
        "Outlook & Future Improvements:",
        "1. NLP Integration:",
        "   - Automatically scrape material properties from research papers.",
        "2. Advanced Decision Making (AHP):",
        "   - Replace simple weighted scoring with Analytic Hierarchy Process.",
        "3. Uncertainty Quantification:",
        "   - Provide error bars for predictions on 'out-of-distribution' geometries.",
        "4. Expansion:",
        "   - Add Fatigue analysis module."
    ]
    slide = add_slide("Conclusion & Future Work", content)
    notes = slide.notes_slide.notes_text_frame
    notes.text = ("Moving forward, there is a lot of potential.\n"
                  "We could use NLP to read material datasheets automatically.\n"
                  "We could implement AHP for better decision making.\n"
                  "And ultimately, adding fatigue analysis would make this a fully industrial-grade tool.")

    # --- SLIDE 20: Thank You ---
    slide = add_slide("Thank You", layout_idx=0)
    title = slide.shapes.title
    title.text = "Thank You for Your Attention"
    subtitle = slide.placeholders[1]
    subtitle.text = "I am now happy to take your questions.\n\nFarah Hany"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Thank you for listening. I am now open to any questions you may have."

    # Save
    prs.save('Bachelor_Thesis_Defense_20min_Farah_Hany.pptx')
    print("Extended presentation saved successfully.")

if __name__ == "__main__":
    create_extended_presentation()